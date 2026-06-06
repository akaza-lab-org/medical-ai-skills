"""
clinical-lecture-ppt-refinement: 監査ユーティリティ

PPTX の現状を把握するための読み取り専用ツール群。
編集前と編集後の両方で実行し、視覚チェックの前段にする。

使い方:
    python audit_slides.py deck.pptx                  # 全スライド概要
    python audit_slides.py deck.pptx --slide 7        # 1枚詳細
    python audit_slides.py deck.pptx --margins        # 下余白チェック
"""
from __future__ import annotations
import sys
import argparse
from pptx import Presentation


def _xy(sh):
    try:
        L = sh.left / 914400 if sh.left is not None else 0
        T = sh.top / 914400 if sh.top is not None else 0
        W = sh.width / 914400 if sh.width is not None else 0
        H = sh.height / 914400 if sh.height is not None else 0
        return L, T, W, H
    except Exception:
        return 0, 0, 0, 0


def summary(prs):
    print(f"Total slides: {len(prs.slides)}")
    print(f"Slide size: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")
    print()
    print(f"{'#':>3}  {'Layout':<24}  Title")
    print("-" * 80)
    for i, slide in enumerate(prs.slides, 1):
        # Prefer the title at the standard position (1.10, 0.55)
        title = None
        for sh in slide.shapes:
            if not sh.has_text_frame: continue
            L, T, _, _ = _xy(sh)
            if abs(L - 1.10) < 0.05 and abs(T - 0.55) < 0.05:
                t = sh.text_frame.text.strip()
                if t:
                    title = t[:60]
                    break
        if title is None:
            for sh in slide.shapes:
                if sh.has_text_frame:
                    t = sh.text_frame.text.strip()
                    if t:
                        title = t.replace('\n', ' / ')[:60]
                        break
        print(f"{i:3d}  {slide.slide_layout.name[:24]:<24}  {title or '(no title)'}")


def detail(prs, n):
    slide = prs.slides[n-1]
    print(f"========== Slide {n} (Layout: {slide.slide_layout.name}) ==========")
    for j, sh in enumerate(slide.shapes):
        L, T, W, H = _xy(sh)
        print(f"  [{j}] ({L:.2f},{T:.2f}) {W:.2f}x{H:.2f}  {type(sh).__name__}")
        if not sh.has_text_frame:
            continue
        for k, para in enumerate(sh.text_frame.paragraphs):
            runs = []
            for run in para.runs:
                f = run.font
                sz = f.size.pt if f.size else None
                b = f.bold
                c = None
                try:
                    if f.color and f.color.rgb:
                        c = str(f.color.rgb)
                except Exception:
                    pass
                runs.append(f"<sz={sz} b={b} c={c}>{run.text!r}</r>")
            print(f"      p{k}: " + (' '.join(runs) if runs else '(empty)'))


def margins(prs):
    print(f"{'#':>3}  {'Bottom y':>10}  {'Slack':>8}  Visual content end?")
    print("-" * 70)
    slide_h = prs.slide_height / 914400
    for i, slide in enumerate(prs.slides, 1):
        max_bottom = 0
        for sh in slide.shapes:
            _, T, _, H = _xy(sh)
            if T < 1.0:  # title region
                continue
            max_bottom = max(max_bottom, T + H)
        slack = slide_h - max_bottom
        flag = '  ←余白多' if slack > 1.5 else ''
        print(f"{i:3d}  {max_bottom:>10.2f}  {slack:>8.2f}{flag}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('pptx', help='Path to .pptx')
    ap.add_argument('--slide', type=int, help='Show detail for one slide (1-based)')
    ap.add_argument('--margins', action='store_true',
                    help='List bottom margins to find rebalance candidates')
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    if args.margins:
        margins(prs)
    elif args.slide:
        detail(prs, args.slide)
    else:
        summary(prs)


if __name__ == '__main__':
    main()
