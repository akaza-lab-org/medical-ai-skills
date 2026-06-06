"""
clinical-lecture-ppt-refinement: 共通ヘルパー

医療講演 PowerPoint デッキを再現性高く整えるための関数群。
python-pptx の API では難しい段落単位の書式制御を、lxml で
直接 OOXML を構築することで安定化させる。

使い方:
    from pptx_helpers import (
        C_HEAD, C_BODY, C_RED,
        make_para, make_mixed_para, replace_paragraphs,
        add_bar, add_title, add_subtitle, delete_shape,
    )

    add_bar(slide)                          # 標準アクセントバー
    add_title(slide, "スライドタイトル")     # 標準タイトル
    replace_paragraphs(body.text_frame, [
        make_para("■  見出し", sz=2200, bold=True, color=C_HEAD),
        make_para("本文行", sz=1700, color=C_BODY, space_before=400),
        make_mixed_para([
            ("通常文 ", 1700, False, C_BODY, False),
            ("赤強調", 1800, True, C_RED, False),
        ], space_before=400),
    ])
"""
from __future__ import annotations
import html
from lxml import etree
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ── Color palette for medical lectures ──────────────────────────
C_TITLE  = "0F172A"   # near-black — slide titles
C_HEAD   = "0284C7"   # blue — ■ section headers
C_BODY   = "334155"   # dark slate — default body text
C_QUOTE  = "475569"   # gray — quoted patient voice
C_SUBTLE = "64748B"   # subtle gray — "before" half of before/after
C_RED    = "C00000"   # red — take-home, abnormal values
C_ORANGE = "B45309"   # amber — caveats, ⚠
C_ACCENT = "29ABE2"   # accent-bar blue

# ── Standard layout positions (inches, for 13.33 x 7.5 deck) ────
BAR_POS     = dict(x=0.80, y=0.60, w=0.10, h=0.70)
TITLE_POS   = dict(x=1.10, y=0.55, w=11.00, h=0.80)
SUBTITLE_POS = dict(x=1.10, y=1.38, w=11.00, h=0.30)
BODY_TOP_Y_WITH_SUBTITLE    = 1.80
BODY_TOP_Y_WITHOUT_SUBTITLE = 1.55

# ── Internal: text escaping ─────────────────────────────────────
def _esc(t: str) -> str:
    return html.escape(t, quote=False)

# ── Paragraph builders ──────────────────────────────────────────
def make_para(
    text: str,
    sz: int,
    *,
    bold: bool = False,
    color: str = C_BODY,
    space_before: int = 0,
    italic: bool = False,
    left_margin_emu: int = 0,
    no_bullet: bool = True,
    line_spc: int | None = None,
    font: str = "Meiryo",
):
    """Build a 1-run paragraph as <a:p> XML.

    Args:
        sz: font size in hundredths of a point (1800 = 18pt)
        space_before: space above paragraph in hundredths of pt (400 = 4pt)
        line_spc: line spacing in per-mille (120000 = 120%)
        left_margin_emu: paragraph left margin in EMU (use Emu(0.35))
        no_bullet: insert <a:buNone/> so placeholder bullets are suppressed
    """
    spc_xml = (
        f'<a:spcBef><a:spcPts val="{space_before}"/></a:spcBef>'
        if space_before else ''
    )
    lns_xml = (
        f'<a:lnSpc><a:spcPct val="{line_spc}"/></a:lnSpc>'
        if line_spc else ''
    )
    bu_xml = '<a:buNone/>' if no_bullet else ''
    margl_attr = f' marL="{left_margin_emu}"' if left_margin_emu else ''
    b = '1' if bold else '0'
    i = '1' if italic else '0'

    return etree.fromstring(f'''<a:p xmlns:a="{NS_A}">
  <a:pPr{margl_attr}>{lns_xml}{spc_xml}{bu_xml}<a:defRPr sz="{sz}" b="{b}" i="{i}">
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:latin typeface="{font}"/>
  </a:defRPr></a:pPr>
  <a:r><a:t>{_esc(text)}</a:t></a:r>
</a:p>''')


def make_mixed_para(
    parts: list[tuple[str, int, bool, str, bool]],
    *,
    space_before: int = 0,
    left_margin_emu: int = 0,
    no_bullet: bool = True,
    line_spc: int | None = None,
    font: str = "Meiryo",
):
    """Paragraph with multiple inline runs of different style.

    Args:
        parts: list of (text, sz, bold, color_hex, italic) tuples
    """
    spc_xml = (
        f'<a:spcBef><a:spcPts val="{space_before}"/></a:spcBef>'
        if space_before else ''
    )
    lns_xml = (
        f'<a:lnSpc><a:spcPct val="{line_spc}"/></a:lnSpc>'
        if line_spc else ''
    )
    bu_xml = '<a:buNone/>' if no_bullet else ''
    margl_attr = f' marL="{left_margin_emu}"' if left_margin_emu else ''

    runs_xml = ""
    for (txt, sz, bold, color, italic) in parts:
        b = '1' if bold else '0'
        i = '1' if italic else '0'
        runs_xml += f'''<a:r>
  <a:rPr sz="{sz}" b="{b}" i="{i}">
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:latin typeface="{font}"/>
  </a:rPr>
  <a:t>{_esc(txt)}</a:t>
</a:r>'''
    return etree.fromstring(f'''<a:p xmlns:a="{NS_A}">
  <a:pPr{margl_attr}>{lns_xml}{spc_xml}{bu_xml}</a:pPr>
  {runs_xml}
</a:p>''')


def replace_paragraphs(text_frame, paragraphs):
    """Replace all <a:p> in text_frame with the new list."""
    tx = text_frame._txBody
    for p in tx.findall(f'{{{NS_A}}}p'):
        tx.remove(p)
    for p in paragraphs:
        tx.append(p)

# ── Shape builders ──────────────────────────────────────────────
def add_bar(slide, x: float = 0.80, y: float = 0.60,
            w: float = 0.10, h: float = 0.70, color: str = C_ACCENT):
    """Add the standard left accent bar at the title row."""
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    return shape


def add_title(slide, text: str, *, x: float = 1.10, y: float = 0.55,
              w: float = 11.0, h: float = 0.80,
              sz: int = 2400, bold: bool = True, color: str = C_TITLE):
    """Add a standard-position title textbox."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    replace_paragraphs(tb.text_frame, [
        make_para(text, sz=sz, bold=bold, color=color),
    ])
    return tb


def add_subtitle(slide, text: str, *, x: float = 1.10, y: float = 1.38,
                 w: float = 11.0, h: float = 0.30,
                 sz: int = 1500, color: str = C_QUOTE):
    """Add a subtitle line just below the title."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    replace_paragraphs(tb.text_frame, [
        make_para(text, sz=sz, color=color),
    ])
    return tb


def add_rect(slide, l: float, t: float, w: float, h: float,
             fill_hex: str = "EFF6FF", line_hex: str | None = None,
             line_width_emu: int = 9525):
    """Add a rectangle. Use for colored boxes in multi-column layouts."""
    shape = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    if line_hex:
        shape.line.color.rgb = RGBColor.from_string(line_hex)
        shape.line.width = Emu(line_width_emu)
    else:
        shape.line.fill.background()
    return shape

# ── Shape position helpers ──────────────────────────────────────
def set_pos(sh, *, x: float | None = None, y: float | None = None,
            w: float | None = None, h: float | None = None):
    """Update shape geometry; any None is left unchanged."""
    if x is not None: sh.left = Inches(x)
    if y is not None: sh.top = Inches(y)
    if w is not None: sh.width = Inches(w)
    if h is not None: sh.height = Inches(h)


def shape_left_in(sh):   return sh.left / 914400 if sh.left is not None else 0
def shape_top_in(sh):    return sh.top / 914400 if sh.top is not None else 0
def shape_width_in(sh):  return sh.width / 914400 if sh.width is not None else 0
def shape_height_in(sh): return sh.height / 914400 if sh.height is not None else 0


def find_shape_at(slide, target_x: float, target_y: float,
                  tol: float = 0.1):
    """Find the first text frame shape near (target_x, target_y) in inches."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        try:
            L = shape_left_in(sh); T = shape_top_in(sh)
        except Exception:
            continue
        if abs(L - target_x) < tol and abs(T - target_y) < tol:
            return sh
    return None


def delete_shape(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)

# ── Slide ordering helpers ──────────────────────────────────────
def insert_slide_after(prs, layout, after_idx_0based: int):
    """Add a slide using `layout` and move it just after `after_idx_0based`."""
    new_slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    new_el = xml_slides[-1]
    xml_slides.remove(new_el)
    xml_slides.insert(after_idx_0based + 1, new_el)
    return new_slide


def delete_slides(prs, idx_1based_list):
    """Delete multiple slides safely, iterating in reverse to keep indices stable."""
    for idx_1based in sorted(set(idx_1based_list), reverse=True):
        idx = idx_1based - 1
        el = prs.slides._sldIdLst[idx]
        rId = el.get(qn('r:id'))
        prs.slides._sldIdLst.remove(el)
        prs.part.drop_rel(rId)


def get_blank_layout(prs):
    """Return a blank slide layout, falling back to layouts[6]."""
    for layout in prs.slide_layouts:
        if 'Blank' in layout.name or layout.name == '白紙':
            return layout
    return prs.slide_layouts[6]

# ── Placeholder → textbox migration ─────────────────────────────
def replace_placeholder_with_textbox(slide, placeholder,
                                     paragraphs,
                                     *, x=None, y=None, w=None, h=None):
    """Delete a SlidePlaceholder (often the source of auto bullets / inherited
    font colors) and add a regular textbox at the same (or overridden) position.
    """
    L = shape_left_in(placeholder) if x is None else x
    T = shape_top_in(placeholder) if y is None else y
    W = shape_width_in(placeholder) if w is None else w
    H = shape_height_in(placeholder) if h is None else h
    delete_shape(slide, placeholder)
    tb = slide.shapes.add_textbox(Inches(L), Inches(T), Inches(W), Inches(H))
    tb.text_frame.word_wrap = True
    replace_paragraphs(tb.text_frame, paragraphs)
    return tb
