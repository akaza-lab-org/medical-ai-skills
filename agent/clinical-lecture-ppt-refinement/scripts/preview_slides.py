"""
clinical-lecture-ppt-refinement: 視覚プレビュー

指定スライドを PowerPoint COM で PNG にエクスポート。
編集前後で同じスライドを書き出し、視覚的に diff を取るのが基本。

使い方:
    python preview_slides.py deck.pptx 7 13 17 18
    python preview_slides.py deck.pptx --all --out C:\\tmp\\preview
"""
from __future__ import annotations
import argparse
import os
import sys

try:
    import win32com.client
except ImportError:
    print("[ERR] pywin32 is required: pip install pywin32", file=sys.stderr)
    sys.exit(1)


def export(pptx_path: str, slide_indices: list[int],
           out_dir: str, w: int = 1920, h: int = 1080):
    os.makedirs(out_dir, exist_ok=True)
    app = win32com.client.Dispatch("PowerPoint.Application")
    # COM in background mode can fail to export; keep visible.
    app.Visible = True
    prs = app.Presentations.Open(pptx_path, ReadOnly=True,
                                  Untitled=False, WithWindow=False)
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    try:
        for n in slide_indices:
            out = os.path.join(out_dir, f"{base}_s{n:02d}.png")
            prs.Slides(n).Export(out, "PNG", w, h)
            print(f"  Slide {n:2d} -> {out}")
    finally:
        prs.Close()
        app.Quit()


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('pptx')
    ap.add_argument('slides', nargs='*', type=int,
                    help='1-based slide indices to export')
    ap.add_argument('--all', action='store_true', help='Export every slide')
    ap.add_argument('--out', default='./preview', help='Output directory')
    ap.add_argument('--width', type=int, default=1920)
    ap.add_argument('--height', type=int, default=1080)
    args = ap.parse_args()

    if args.all:
        # Use python-pptx just to count
        from pptx import Presentation
        n = len(Presentation(args.pptx).slides)
        indices = list(range(1, n + 1))
    else:
        if not args.slides:
            ap.error('Provide slide numbers or --all')
        indices = args.slides

    pptx = os.path.abspath(args.pptx)
    out  = os.path.abspath(args.out)
    export(pptx, indices, out, args.width, args.height)


if __name__ == '__main__':
    main()
