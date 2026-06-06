import os
import pdfplumber
from pptx import Presentation
import json
import argparse

def get_pdf_titles(path):
    titles = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    first_line = text.split('\n')[0].strip()
                    titles.append({"page": i + 1, "title": first_line})
                else:
                    titles.append({"page": i + 1, "title": "[Image/No Text]"})
    except Exception as e:
        return str(e)
    return titles

def get_pptx_titles(path):
    titles = []
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            title = "[No Title]"
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
            elif len(slide.shapes) > 0:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        title = shape.text.strip().split('\n')[0]
                        break
            titles.append({"page": i + 1, "title": title})
    except Exception as e:
        return str(e)
    return titles

def main():
    parser = argparse.ArgumentParser(description='Extract titles and page numbers from PPTX/PDF files.')
    parser.add_argument('files', nargs='+', help='List of files to process')
    args = parser.parse_args()

    results = {}
    for path in args.files:
        if not os.path.exists(path):
            continue
        key = os.path.basename(path)
        if path.endswith(".pptx"):
            results[key] = get_pptx_titles(path)
        elif path.endswith(".pdf"):
            results[key] = get_pdf_titles(path)
    
    print(json.dumps(results, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()
